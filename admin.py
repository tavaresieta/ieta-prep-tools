"""
IETA Admin — Gestão da Base de Conhecimento
Página separada para:
  1. Upload de novos documentos (PDF/DOCX/PPTX) com extração automática
  2. Validação de keywords sugeridas pelo Claude antes de salvar
  3. Busca por texto livre nos documentos (Caminho B)
 
Como usar:
  streamlit run admin.py
"""
 
import streamlit as st
from pathlib import Path
from PyPDF2 import PdfReader
from docx import Document
import json
import os
import anthropic
from datetime import datetime
 
# ── Tenta importar pptx, mas não obriga ──────────────────────────────────────
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
 
# ─────────────────────────────────────────────────────────────────────────────
# CONFIG
# ─────────────────────────────────────────────────────────────────────────────
DOCS_FOLDER   = Path("documents")
META_FOLDER   = Path("metadata")
META_FILE     = META_FOLDER / "keywords_metadata.json"
INDEX_FILE    = META_FOLDER / "keywords_metadata_index.json"
 
DOC_TYPES = [
    "IETA Report",
    "IETA Position / Discussion Paper",
    "IETA Presentation / Slide Deck",
    "External Report",
    "External Paper / Research",
    "External Presentation / Slide Deck",
]
 
# Vocabulário padrão — usado apenas se vocabulary.json ainda não existir
_DEFAULT_VOCABULARY = [
    # Regulação
    "artigo-6.2","artigo-6.4","sbce","cdm-transition","compliance-market",
    "corresponding-adjustments","mrv","transparency","baseline",
    # Mercado
    "ecora","vcm","voluntary-carbon-market","certificação","credito-carbono",
    "open-coalition","financiamento","tfff","ecoinvest","fundo-clima",
    "jredd","cadtrust","registry","verificação",
    # Setores
    "nbs","nature-based-solutions","floresta","reflorestamento",
    "biocombustiveis","etanol","biodiesel","saf",
    "energia","renovaveis","solar","eolica","hidro",
    "cimento","mineracao","aviacao","hard-to-abate",
    "agricultura","pecuaria","uso-do-solo",
    # Conceitos
    "integridade","adicionalidade","permanencia","metodologia",
    "escopo-1","escopo-2","escopo-3","inventario",
    "net-zero","neutralidade","compensacao","remocao",
    "co-beneficios","sdgs","comunidades-tradicionais",
    # Geografia
    "brasil","amazonia","cerrado","mata-atlantica",
    "america-latina","global","bilateral",
    # Atores
    "governo","empresa-privada","ong","academia",
    "comunidade-local","indigena","unfccc",
]
 
@st.cache_data
def get_vocabulary() -> list:
    """Sempre lê o vocabulary.json se existir — fonte única da verdade.
    Decorada com cache_data para ser invalidada junto com st.cache_data.clear()."""
    vocab_file = META_FOLDER / "vocabulary.json"
    if vocab_file.exists():
        try:
            return json.loads(vocab_file.read_text(encoding="utf-8"))
        except Exception:
            pass
    return _DEFAULT_VOCABULARY
 
# Usado como fallback em contextos que ainda referenciam KEYWORD_VOCABULARY
KEYWORD_VOCABULARY = _DEFAULT_VOCABULARY
 
st.set_page_config(
    page_title="IETA Admin",
    page_icon="⚙️",
    layout="wide",
)
 
# ─────────────────────────────────────────────────────────────────────────────
# HELPERS — IO
# ─────────────────────────────────────────────────────────────────────────────
 
def ensure_folders():
    DOCS_FOLDER.mkdir(parents=True, exist_ok=True)
    META_FOLDER.mkdir(parents=True, exist_ok=True)
 
def load_metadata() -> dict:
    if META_FILE.exists():
        return json.loads(META_FILE.read_text(encoding="utf-8"))
    return {}
 
def save_metadata(meta: dict):
    META_FILE.write_text(json.dumps(meta, indent=2, ensure_ascii=False), encoding="utf-8")
    # rebuild index
    index: dict = {}
    for doc_id, doc in meta.items():
        for kw in doc.get("keywords", []):
            index.setdefault(kw, []).append(doc_id)
    INDEX_FILE.write_text(json.dumps(index, indent=2, ensure_ascii=False), encoding="utf-8")
 
# ─────────────────────────────────────────────────────────────────────────────
# HELPERS — EXTRAÇÃO DE TEXTO
# ─────────────────────────────────────────────────────────────────────────────
 
def extract_text(uploaded_file) -> str:
    """Extrai texto de PDF, DOCX ou PPTX a partir de um UploadedFile."""
    name = uploaded_file.name.lower()
    raw  = uploaded_file.read()
 
    if name.endswith(".pdf"):
        from io import BytesIO
        try:
            reader = PdfReader(BytesIO(raw))
            if reader.is_encrypted:
                st.warning(
                    f"⚠️ `{uploaded_file.name}` está criptografado e não pode ser lido. "
                    "Salve uma versão sem senha e faça upload novamente."
                )
                return ""
            return "\n".join(p.extract_text() or "" for p in reader.pages)
        except Exception as e:
            st.warning(f"⚠️ Não foi possível ler `{uploaded_file.name}`: {e}")
            return ""
 
    if name.endswith(".docx"):
        from io import BytesIO
        doc = Document(BytesIO(raw))
        return "\n".join(p.text for p in doc.paragraphs)
 
    if name.endswith(".pptx"):
        if not HAS_PPTX:
            st.warning("python-pptx não instalado — instale com `pip install python-pptx`")
            return ""
        from io import BytesIO
        prs  = Presentation(BytesIO(raw))
        lines = []
        for i, slide in enumerate(prs.slides, 1):
            lines.append(f"\n--- SLIDE {i} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    lines.append(shape.text)
        return "\n".join(lines)
 
    return ""
 
# ─────────────────────────────────────────────────────────────────────────────
# HELPERS — KEYWORDS VIA CLAUDE
# ─────────────────────────────────────────────────────────────────────────────
 
def suggest_keywords_claude(text: str, n: int = 10) -> list[str]:
    """Pede ao Claude keywords do vocabulário + até 5 novas sugestões livres."""
    api_key = os.environ.get("ANTHROPIC_API_KEY") or st.secrets.get("ANTHROPIC_API_KEY", "")
    if not api_key:
        st.error("⚠️ ANTHROPIC_API_KEY não configurada.")
        return []
 
    client = anthropic.Anthropic(api_key=api_key)
    # usa vocabulário customizado se existir
    vocab_file = META_FOLDER / "vocabulary.json"
    vocab = get_vocabulary()
    vocab_str = ", ".join(vocab)
 
    prompt = f"""Você é um especialista em mercados de carbono e clima.
Analise o texto abaixo e retorne EXATAMENTE duas seções em JSON:
 
{{
  "from_vocabulary": ["até {n} keywords do vocabulário que se aplicam ao texto"],
  "new_suggestions": ["até 5 keywords novas, relevantes, que NÃO estão no vocabulário — em formato kebab-case"]
}}
 
VOCABULÁRIO DISPONÍVEL:
{vocab_str}
 
TEXTO (primeiros 5000 caracteres):
{text[:5000]}
 
Responda APENAS com o JSON, sem markdown, sem explicações."""
 
    try:
        msg = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=400,
            messages=[{"role": "user", "content": prompt}],
        )
        raw = msg.content[0].text.strip()
        # remove possíveis backticks
        raw = raw.replace("```json", "").replace("```", "").strip()
        data = json.loads(raw)
        return data
    except Exception as e:
        st.error(f"Erro na API Claude: {e}")
        return {"from_vocabulary": [], "new_suggestions": []}
 
# ─────────────────────────────────────────────────────────────────────────────
# HELPERS — BUSCA POR TEXTO LIVRE
# ─────────────────────────────────────────────────────────────────────────────
 
def search_documents(query: str, documents: dict) -> list[dict]:
    """
    Busca simples por texto livre:
    - divide query em termos
    - conta ocorrências em cada documento
    - retorna lista ordenada por score, com trechos contextuais
    """
    if not query.strip():
        return []
 
    terms = [t.lower().strip() for t in query.split() if len(t.strip()) > 2]
    if not terms:
        return []
 
    results = []
    for doc_id, doc in documents.items():
        content_lower = doc["full_content"].lower()
        score = sum(content_lower.count(t) for t in terms)
        if score == 0:
            continue
 
        # Encontra trecho contextual do primeiro termo
        snippet = ""
        for term in terms:
            idx = content_lower.find(term)
            if idx != -1:
                start = max(0, idx - 120)
                end   = min(len(doc["full_content"]), idx + 200)
                snippet = "…" + doc["full_content"][start:end].replace("\n", " ") + "…"
                break
 
        results.append({
            "doc_id":   doc_id,
            "filename": doc["filename"],
            "score":    score,
            "snippet":  snippet,
            "size_kb":  doc["size_kb"],
        })
 
    results.sort(key=lambda x: x["score"], reverse=True)
    return results
 
# ─────────────────────────────────────────────────────────────────────────────
# CARREGA DOCUMENTOS TXT
# ─────────────────────────────────────────────────────────────────────────────
 
@st.cache_data
def load_documents():
    docs = {}
    for fp in DOCS_FOLDER.glob("*.txt"):
        try:
            content = fp.read_text(encoding="utf-8")
            if content.strip():
                docs[fp.stem] = {
                    "filename":     fp.name,
                    "full_content": content,
                    "size_kb":      round(len(content) / 1024, 1),
                    "char_count":   len(content),
                }
        except Exception:
            pass
    return docs
 
# ─────────────────────────────────────────────────────────────────────────────
# UI
# ─────────────────────────────────────────────────────────────────────────────
 
ensure_folders()
 
st.markdown("""
<style>
  /* Fonte principal */
  @import url('https://fonts.googleapis.com/css2?family=IBM+Plex+Mono:wght@400;600&family=IBM+Plex+Sans:wght@300;400;600&display=swap');
  html, body, [class*="css"] { font-family: 'IBM Plex Sans', sans-serif; }
  h1, h2, h3 { font-family: 'IBM Plex Mono', monospace; letter-spacing: -0.03em; }
 
  /* Tabs */
  .stTabs [data-baseweb="tab-list"] { gap: 0; border-bottom: 2px solid #1a1a2e; }
  .stTabs [data-baseweb="tab"] {
    font-family: 'IBM Plex Mono', monospace;
    font-size: 0.8rem;
    font-weight: 600;
    letter-spacing: 0.05em;
    text-transform: uppercase;
    padding: 0.6rem 1.4rem;
    border-radius: 0;
    border: 2px solid transparent;
    border-bottom: none;
    color: #666;
  }
  .stTabs [aria-selected="true"] {
    background: #1a1a2e !important;
    color: #e8f4f8 !important;
    border-color: #1a1a2e !important;
  }
 
  /* Cards de resultado */
  .result-card {
    background: #f7f9fc;
    border-left: 4px solid #1a1a2e;
    padding: 1rem 1.2rem;
    margin: 0.6rem 0;
    border-radius: 0 6px 6px 0;
  }
  .result-card .score-badge {
    display: inline-block;
    background: #1a1a2e;
    color: #7effd4;
    font-family: 'IBM Plex Mono', monospace;
    font-size: 0.72rem;
    padding: 2px 8px;
    border-radius: 2px;
    margin-right: 8px;
  }
  .result-card .snippet { color: #444; font-size: 0.88rem; margin-top: 0.4rem; line-height: 1.5; }
  .kw-chip {
    display: inline-block;
    background: #e8f4f8;
    color: #1a1a2e;
    font-family: 'IBM Plex Mono', monospace;
    font-size: 0.72rem;
    padding: 3px 10px;
    border-radius: 20px;
    margin: 2px;
    border: 1px solid #b8d8e8;
  }
  .kw-chip-new {
    background: #fff8e1;
    border-color: #f0c040;
    color: #7a5c00;
  }
</style>
""", unsafe_allow_html=True)
 
# Header
col_h1, col_h2 = st.columns([3, 1])
with col_h1:
    st.markdown("# ⚙️ IETA Admin")
    st.caption("Gestão da base de conhecimento · Upload · Keywords · Busca")
with col_h2:
    if st.button("🔄 Recarregar cache", use_container_width=True):
        st.cache_data.clear()
        st.success("Cache limpo!")
        st.rerun()
 
st.markdown("---")
 
# ── Quatro abas principais ────────────────────────────────────────────────────
tab_overview, tab_upload, tab_keywords, tab_search = st.tabs([
    "📋  Visão Geral",
    "📤  Upload & Processamento",
    "🏷️  Keywords & Validação",
    "🔍  Busca por Texto",
])
 
# =============================================================================
# ABA 0 — VISÃO GERAL
# =============================================================================
with tab_overview:
    st.markdown("### Visão Geral da Base")
 
    metadata = load_metadata()
 
    # ── Sincronizar docs sem metadata ─────────────────────────────────────────
    docs_no_meta = [
        fp for fp in DOCS_FOLDER.glob("*.txt")
        if fp.stem not in metadata
    ]
    if docs_no_meta:
        st.warning(f"⚠️ {len(docs_no_meta)} documento(s) sem metadata encontrado(s):")
        for fp in docs_no_meta:
            st.caption(f"  · `{fp.name}`")
        if st.button("➕ Registrar esses documentos (keywords vazias)", type="primary"):
            for fp in docs_no_meta:
                text = fp.read_text(encoding="utf-8")
                metadata[fp.stem] = {
                    "filename":       fp.name,
                    "keywords":       [],
                    "doc_type":       "",
                    "char_count":     len(text),
                    "word_count":     len(text.split()),
                    "processed_date": datetime.now().isoformat(),
                }
            save_metadata(metadata)
            st.cache_data.clear()
            st.success(f"✅ {len(docs_no_meta)} documento(s) registrado(s)! Atribua as keywords abaixo.")
            st.rerun()
        st.markdown("---")
 
    col_esq, col_dir = st.columns([1, 1], gap="large")
 
    # ── Coluna esquerda: keywords por documento ───────────────────────────────
    with col_esq:
        st.markdown("#### Keywords por documento")
        st.caption("Edite diretamente — uma keyword por linha. Salve ao terminar.")
 
        if metadata:
            # Monta texto editável: "nome_arquivo: kw1, kw2, kw3"
            linhas = []
            for doc_id, doc in sorted(metadata.items(), key=lambda x: x[1].get("filename", x[0])):
                kws = ", ".join(doc.get("keywords", []))
                linhas.append(f"{doc.get('filename', doc_id)}: {kws}")
            texto_inicial = "\n".join(linhas)
 
            texto_editado = st.text_area(
                "Um documento por linha  —  `nome_arquivo.txt: kw1, kw2, kw3`",
                value=texto_inicial,
                height=480,
                key="overview_docs",
            )
 
            if st.button("💾 Salvar alterações dos documentos", use_container_width=True, type="primary"):
                erros = []
                # Constrói mapa filename -> doc_id
                fname_to_id = {v.get("filename", k): k for k, v in metadata.items()}
 
                for linha in texto_editado.strip().splitlines():
                    linha = linha.strip()
                    if not linha:
                        continue
                    if ":" not in linha:
                        erros.append(f"Linha ignorada (sem ':'): {linha}")
                        continue
                    fname, _, kws_raw = linha.partition(":")
                    fname = fname.strip()
                    kws = [k.strip() for k in kws_raw.split(",") if k.strip()]
 
                    doc_id = fname_to_id.get(fname)
                    if not doc_id:
                        # tenta pelo stem
                        doc_id = fname_to_id.get(fname + ".txt") or Path(fname).stem
                        if doc_id not in metadata:
                            erros.append(f"Documento não encontrado: {fname}")
                            continue
 
                    metadata[doc_id]["keywords"] = kws
                    metadata[doc_id]["processed_date"] = datetime.now().isoformat()
 
                save_metadata(metadata)
                st.cache_data.clear()
                if erros:
                    for e in erros:
                        st.warning(e)
                st.success("✅ Keywords salvas!")
                st.rerun()
        else:
            st.info("Nenhum documento com metadata ainda.")
 
    # ── Coluna direita: vocabulário padrão ───────────────────────────────────
    with col_dir:
        st.markdown("#### Vocabulário padrão")
        st.caption("Edite a lista mestre de keywords. Uma por linha.")
 
        # Carrega vocabulário customizado se existir, senão usa o padrão do código
        vocab_file = META_FOLDER / "vocabulary.json"
        if vocab_file.exists():
            vocab_atual = json.loads(vocab_file.read_text(encoding="utf-8"))
        else:
            vocab_atual = get_vocabulary()
 
        vocab_texto = "\n".join(sorted(vocab_atual))
 
        vocab_editado = st.text_area(
            "Uma keyword por linha",
            value=vocab_texto,
            height=480,
            key="overview_vocab",
        )
 
        if st.button("💾 Salvar vocabulário", use_container_width=True):
            novas = [k.strip() for k in vocab_editado.splitlines() if k.strip()]
            novas = sorted(set(novas))  # deduplica e ordena
            vocab_file.write_text(json.dumps(novas, ensure_ascii=False, indent=2), encoding="utf-8")
            st.cache_data.clear()  # invalida get_vocabulary() e todos os outros caches
            st.success(f"✅ Vocabulário salvo com {len(novas)} keywords!")
            st.rerun()
 
        st.markdown("---")
        # Estatísticas rápidas
        from collections import Counter
        todas_kws = []
        for d in metadata.values():
            todas_kws.extend(d.get("keywords", []))
        contagem = Counter(todas_kws)
 
        st.markdown("#### Keywords mais usadas na base")
        if contagem:
            for kw, n in contagem.most_common(15):
                fora = "⚠️" if kw not in vocab_atual else ""
                st.markdown(f"`{kw}` — {n} doc(s) {fora}")
            if any(kw not in vocab_atual for kw in contagem):
                st.caption("⚠️ = keyword usada em documentos mas não está no vocabulário padrão")
        else:
            st.info("Nenhuma keyword atribuída ainda.")
 
# =============================================================================
# ABA 1 — UPLOAD
# =============================================================================
with tab_upload:
    st.markdown("### Adicionar novos documentos")
    st.markdown(
        "Faça upload de PDFs, DOCXs ou PPTXs. O texto será extraído, "
        "salvo em `documents/` e o Claude vai sugerir keywords para validação."
    )
 
    uploaded_files = st.file_uploader(
        "Arraste ou selecione arquivos",
        type=["pdf", "docx", "pptx"],
        accept_multiple_files=True,
        label_visibility="collapsed",
    )
 
    if uploaded_files:
        st.markdown(f"**{len(uploaded_files)} arquivo(s) selecionado(s):**")
 
        metadata = load_metadata()
        to_process = []
 
        for uf in uploaded_files:
            stem = Path(uf.name).stem
            already = (DOCS_FOLDER / f"{stem}.txt").exists()
            status = "⚠️ já existe" if already else "🆕 novo"
            st.markdown(f"- `{uf.name}` — {status}")
            if not already:
                to_process.append(uf)
 
        if not to_process:
            st.info("Todos os arquivos já estão na base. Use a aba Keywords para rever.")
        else:
            st.markdown("---")
            col_a, col_b = st.columns(2)
            with col_a:
                extract_kw = st.checkbox("Sugerir keywords automaticamente (Claude)", value=True)
            with col_b:
                go_to_validate = st.checkbox("Ir direto para validação após upload", value=True)
 
            doc_type_upload = st.selectbox(
                "**Tipo de documento:**",
                options=DOC_TYPES,
                key="doc_type_upload",
            )
 
            if st.button("🚀 Processar e salvar", type="primary", use_container_width=True):
                new_pending = {}  # doc_id -> {text, suggested}
 
                progress = st.progress(0)
                for i, uf in enumerate(to_process):
                    stem = Path(uf.name).stem
                    st.write(f"📄 Extraindo `{uf.name}`…")
 
                    text = extract_text(uf)
                    if len(text.strip()) < 50:
                        st.warning(f"  ⚠️ Pouco texto extraído de {uf.name} — verifique o arquivo.")
                        continue
 
                    # Salva TXT
                    txt_path = DOCS_FOLDER / f"{stem}.txt"
                    txt_path.write_text(text, encoding="utf-8")
                    st.success(f"  ✅ Salvo: `{txt_path.name}` ({len(text):,} chars)")
 
                    # Sugere keywords
                    suggested = {"from_vocabulary": [], "new_suggestions": []}
                    if extract_kw:
                        with st.spinner(f"  🤖 Claude analisando keywords de {uf.name}…"):
                            suggested = suggest_keywords_claude(text)
 
                    # Armazena em session_state para validação
                    new_pending[stem] = {
                        "filename":    uf.name,
                        "text":        text,
                        "suggested":   suggested,
                        "char_count":  len(text),
                        "word_count":  len(text.split()),
                        "doc_type":    doc_type_upload,
                    }
 
                    progress.progress((i + 1) / len(to_process))
 
                # Guarda pendentes para a aba de validação
                existing = st.session_state.get("pending_validation", {})
                existing.update(new_pending)
                st.session_state["pending_validation"] = existing
                st.cache_data.clear()
 
                st.success(f"✅ {len(new_pending)} documento(s) processado(s)!")
                if go_to_validate and new_pending:
                    st.info("👉 Acesse a aba **Keywords & Validação** para aprovar as keywords sugeridas.")
 
    else:
        # Mostra estatísticas atuais
        st.markdown("#### Base atual")
        docs = load_documents()
        metadata = load_metadata()
 
        if docs:
            total_kb = sum(d["size_kb"] for d in docs.values())
            col1, col2, col3 = st.columns(3)
            col1.metric("Documentos", len(docs))
            col2.metric("Tamanho total", f"{total_kb:.1f} KB")
            col3.metric("Com keywords", sum(1 for d in metadata.values() if d.get("keywords")))
 
            with st.expander(f"Ver todos os {len(docs)} documentos"):
                for doc_id, doc in sorted(docs.items()):
                    kws = metadata.get(doc_id, {}).get("keywords", [])
                    kw_html = "".join(f'<span class="kw-chip">{k}</span>' for k in kws) if kws else "<em style='color:#aaa'>sem keywords</em>"
                    st.markdown(
                        f"**{doc['filename']}** — {doc['char_count']:,} chars<br>{kw_html}",
                        unsafe_allow_html=True,
                    )
        else:
            st.warning("Nenhum documento na base ainda. Faça upload acima!")
 
# =============================================================================
# ABA 2 — KEYWORDS & VALIDAÇÃO
# =============================================================================
with tab_keywords:
    st.markdown("### Validação e edição de keywords")
 
    metadata = load_metadata()
    pending  = st.session_state.get("pending_validation", {})
 
    # ── Documentos aguardando validação ──────────────────────────────────────
    if pending:
        st.markdown(f"#### 🟡 {len(pending)} documento(s) aguardando validação")
 
        approved_all = {}
 
        for doc_id, info in list(pending.items()):
            st.markdown(f"---\n##### 📄 `{info['filename']}`")
            st.caption(f"{info['char_count']:,} caracteres · {info['word_count']:,} palavras")
 
            suggested_vocab = info["suggested"].get("from_vocabulary", [])
            suggested_new   = info["suggested"].get("new_suggestions", [])
 
            doc_type_val = st.selectbox(
                "**Tipo de documento:**",
                options=DOC_TYPES,
                index=DOC_TYPES.index(info.get("doc_type", DOC_TYPES[0])) if info.get("doc_type") in DOC_TYPES else 0,
                key=f"dtype_{doc_id}",
            )
 
            col_left, col_right = st.columns([3, 2])
 
            with col_left:
                st.markdown("**Keywords do vocabulário sugeridas pelo Claude:**")
                if suggested_vocab:
                    approved = st.multiselect(
                        "Selecione as que quer manter (pode remover ou adicionar):",
                        options=get_vocabulary(),
                        default=suggested_vocab,
                        key=f"vocab_{doc_id}",
                    )
                else:
                    approved = st.multiselect(
                        "Nenhuma sugestão — selecione manualmente:",
                        options=get_vocabulary(),
                        default=[],
                        key=f"vocab_{doc_id}",
                    )
 
                # Novas keywords sugeridas fora do vocabulário
                if suggested_new:
                    st.markdown("**Novas keywords sugeridas (fora do vocabulário):**")
                    new_chips = "".join(f'<span class="kw-chip kw-chip-new">✨ {k}</span>' for k in suggested_new)
                    st.markdown(new_chips, unsafe_allow_html=True)
 
                    add_new = st.multiselect(
                        "Quais adicionar ao vocabulário deste documento?",
                        options=suggested_new,
                        default=suggested_new,
                        key=f"new_{doc_id}",
                    )
                else:
                    add_new = []
 
                # Campo para adicionar keywords livres manualmente
                extra_raw = st.text_input(
                    "Adicionar keyword manual (kebab-case, pressione Enter):",
                    key=f"extra_{doc_id}",
                    placeholder="ex: mercado-voluntario-2030",
                )
                extra = [k.strip() for k in extra_raw.split(",") if k.strip()] if extra_raw else []
 
                final_keywords = list(dict.fromkeys(approved + add_new + extra))  # deduplica mantendo ordem
                approved_all[doc_id] = {"info": info, "keywords": final_keywords}
 
            with col_right:
                st.markdown("**Preview do documento:**")
                preview = info["text"][:800].replace("\n", "  \n")
                st.text_area("", value=preview, height=220, disabled=True, key=f"prev_{doc_id}")
 
            st.markdown(
                "**Keywords finais:** " + "".join(f'<span class="kw-chip">{k}</span>' for k in final_keywords),
                unsafe_allow_html=True,
            )
 
        st.markdown("---")
        if st.button("✅ Salvar todas as keywords validadas", type="primary", use_container_width=True):
            for doc_id, data in approved_all.items():
                info = data["info"]
                metadata[doc_id] = {
                    "filename":       info["filename"],
                    "keywords":       data["keywords"],
                    "doc_type":       st.session_state.get(f"dtype_{doc_id}", info.get("doc_type", "")),
                    "char_count":     info["char_count"],
                    "word_count":     info["word_count"],
                    "processed_date": datetime.now().isoformat(),
                }
            save_metadata(metadata)
            st.session_state["pending_validation"] = {}
            st.cache_data.clear()
            st.success(f"✅ Keywords de {len(approved_all)} documento(s) salvas!")
            st.rerun()
 
    else:
        st.info("Nenhum documento aguardando validação. Faça upload na aba anterior.")
 
    # ── Editar keywords de documentos já salvos ───────────────────────────────
    st.markdown("---")
    st.markdown("#### 📚 Editar keywords de documentos existentes")
 
    if metadata:
        doc_options = {v["filename"]: k for k, v in metadata.items()}
        chosen_name = st.selectbox("Selecione um documento para editar:", list(doc_options.keys()))
 
        if chosen_name:
            chosen_id    = doc_options[chosen_name]
            current_kws  = metadata[chosen_id].get("keywords", [])
            current_type = metadata[chosen_id].get("doc_type", DOC_TYPES[0])
 
            doc_type_edit = st.selectbox(
                "**Tipo de documento:**",
                options=DOC_TYPES,
                index=DOC_TYPES.index(current_type) if current_type in DOC_TYPES else 0,
                key="doc_type_edit",
            )
 
            col_e1, col_e2 = st.columns([2, 1])
            with col_e1:
                edited = st.multiselect(
                    "Keywords atuais (edite à vontade):",
                    options=get_vocabulary(),
                    default=[k for k in current_kws if k in get_vocabulary()],
                    key="edit_existing",
                )
                extra_edit_raw = st.text_input(
                    "Keywords customizadas (separadas por vírgula):",
                    value=", ".join(k for k in current_kws if k not in get_vocabulary()),
                    key="edit_custom",
                )
                extra_edit = [k.strip() for k in extra_edit_raw.split(",") if k.strip()]
                final_edit = list(dict.fromkeys(edited + extra_edit))
 
            with col_e2:
                st.markdown("**Keywords atuais:**")
                chips = "".join(f'<span class="kw-chip">{k}</span>' for k in current_kws) or "<em>nenhuma</em>"
                st.markdown(chips, unsafe_allow_html=True)
 
                st.markdown("**Após edição:**")
                chips2 = "".join(f'<span class="kw-chip">{k}</span>' for k in final_edit) or "<em>nenhuma</em>"
                st.markdown(chips2, unsafe_allow_html=True)
 
            if st.button("💾 Salvar edição", use_container_width=True):
                metadata[chosen_id]["keywords"] = final_edit
                metadata[chosen_id]["doc_type"] = doc_type_edit
                metadata[chosen_id]["processed_date"] = datetime.now().isoformat()
                save_metadata(metadata)
                st.cache_data.clear()
                st.success(f"✅ Keywords de `{chosen_name}` atualizadas!")
                st.rerun()
 
            # Botão para regenerar keywords com Claude
            st.markdown("---")
            if st.button("🤖 Re-sugerir keywords com Claude", use_container_width=True):
                txt_path = DOCS_FOLDER / f"{chosen_id}.txt"
                if txt_path.exists():
                    text = txt_path.read_text(encoding="utf-8")
                    with st.spinner("Claude analisando…"):
                        result = suggest_keywords_claude(text)
                    st.markdown("**Sugestão do vocabulário:**")
                    st.write(result.get("from_vocabulary", []))
                    st.markdown("**Novas sugestões:**")
                    st.write(result.get("new_suggestions", []))
                    st.info("Copie as que quiser para o campo acima e salve.")
                else:
                    st.error("Arquivo TXT não encontrado.")
    else:
        st.info("Nenhum documento com metadata. Faça upload e valide keywords primeiro.")
 
# =============================================================================
# ABA 3 — BUSCA POR TEXTO LIVRE
# =============================================================================
with tab_search:
    st.markdown("### Busca por texto livre")
    st.markdown("Pesquise por qualquer termo diretamente no conteúdo dos documentos.")
 
    documents = load_documents()
 
    if not documents:
        st.warning("Base vazia. Faça upload de documentos primeiro.")
    else:
        query = st.text_input(
            "Digite sua busca:",
            placeholder="ex: mercado de carbono biodiesel correspondências",
            label_visibility="collapsed",
        )
 
        col_f1, col_f2 = st.columns([3, 1])
        with col_f2:
            max_results = st.selectbox("Máx. resultados:", [5, 10, 20, 50], index=1)
 
        if query:
            results = search_documents(query, documents)[:max_results]
 
            if results:
                st.markdown(f"**{len(results)} documento(s) encontrado(s)** para `{query}`")
                metadata = load_metadata()
 
                for r in results:
                    doc_id = r["doc_id"]
                    kws    = metadata.get(doc_id, {}).get("keywords", [])
                    kw_html = "".join(f'<span class="kw-chip">{k}</span>' for k in kws)
 
                    st.markdown(f"""
<div class="result-card">
  <span class="score-badge">{r['score']} ocorrências</span>
  <strong>{r['filename']}</strong> <span style="color:#888;font-size:0.8rem">— {r['size_kb']} KB</span>
  <div style="margin-top:6px">{kw_html}</div>
  <div class="snippet">{r['snippet']}</div>
</div>
""", unsafe_allow_html=True)
 
                    with st.expander(f"Ver conteúdo completo de {r['filename']}"):
                        content = documents[doc_id]["full_content"]
                        # Destaca termos buscados
                        highlighted = content
                        st.text_area("", value=highlighted, height=300, disabled=True, key=f"full_{doc_id}")
            else:
                st.warning(f"Nenhum documento contém os termos de `{query}`.")
                st.caption("Dica: tente termos mais curtos ou em inglês se os documentos forem bilíngues.")
 
        else:
            # Mostra overview da base quando não há query
            st.markdown(f"**{len(documents)} documentos na base** — {sum(d['size_kb'] for d in documents.values()):.0f} KB total")
 
            metadata = load_metadata()
            # Nuvem de keywords
            from collections import Counter
            all_kws = []
            for d in metadata.values():
                all_kws.extend(d.get("keywords", []))
            kw_counts = Counter(all_kws).most_common(30)
 
            if kw_counts:
                st.markdown("#### Keywords mais frequentes na base")
                chips = ""
                for kw, count in kw_counts:
                    size = 0.75 + (count / max(c for _, c in kw_counts)) * 0.5
                    chips += f'<span class="kw-chip" style="font-size:{size:.2f}rem">{kw} <strong>({count})</strong></span>'
                st.markdown(chips, unsafe_allow_html=True)