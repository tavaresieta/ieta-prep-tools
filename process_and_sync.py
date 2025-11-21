"""
Workflow Integrado - IETA Knowledge Base
Processa PDFs/DOCs/PPTXs → Chunks se necessário → TXTs → Keywords → Sync GitHub

USO:
    python process_and_sync.py                    # Processa tudo
    python process_and_sync.py --keywords-only    # Só atualiza keywords
    python process_and_sync.py --no-push          # Não faz push no GitHub
    python process_and_sync.py --no-chunk         # Desabilita chunking automático
"""

import os
from pathlib import Path
from datetime import datetime
import json
import subprocess
from typing import Dict, List, Optional
import anthropic

# Seus imports existentes de extração
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

# ============================================================================
# CONFIGURAÇÃO
# ============================================================================

CONFIG = {
    # Pastas locais
    "input_folder": r"C:\Users\maryb\OneDrive\Desktop\IETA_Tests\IETA_Bot\input",
    "documents_folder": r"C:\Users\maryb\OneDrive\Desktop\IETA_Tests\IETA_Bot\documents",
    "metadata_folder": r"C:\Users\maryb\OneDrive\Desktop\IETA_Tests\IETA_Bot\metadata",
    
    # API Keys
    "anthropic_api_key": "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",

    # GitHub
    "github_repo": r"C:\Users\maryb\OneDrive\Desktop\IETA_Tests\IETA_Bot",
    "auto_commit": True,
    "auto_push": True,
    
    # Chunking (NOVO!)
    "auto_chunk": True,
    "chunk_threshold_pages": 20,      # PDFs com mais de 20 páginas = chunk
    "chunk_threshold_chars": 100000,   # TXTs com mais de 100k chars = chunk
    "pages_per_chunk": 15,
    "paragraphs_per_chunk": 100,
    
    # Keywords
    "extract_keywords": True,
    "keywords_per_doc": 5,
    "use_claude_api": True,
    
    # Vocabulário de keywords
    "keyword_vocabulary": [
        # Regulação
        "artigo-6.2", "artigo-6.4", "sbce", "cdm-transition", "compliance-market",
        "corresponding-adjustments", "mrv", "transparency", "baseline",
        
        # Mercado
        "ecora", "vcm", "voluntary-carbon-market", "certificação", "credito-carbono",
        "open-coalition", "financiamento", "tfff", "ecoinvest", "fundo-clima",
        "jredd", "cadtrust", "registry", "verificação",
        
        # Setores
        "nbs", "nature-based-solutions", "floresta", "reflorestamento",
        "biocombustiveis", "etanol", "biodiesel", "saf",
        "energia", "renovaveis", "solar", "eolica", "hidro",
        "cimento", "mineracao", "aviacao", "hard-to-abate",
        "agricultura", "pecuaria", "uso-do-solo",
        
        # Conceitos
        "integridade", "adicionalidade", "permanencia", "metodologia",
        "escopo-1", "escopo-2", "escopo-3", "inventario",
        "net-zero", "neutralidade", "compensacao", "remocao",
        "co-beneficios", "sdgs", "comunidades-tradicionais",
        
        # Geografia
        "brasil", "amazonia", "cerrado", "mata-atlantica",
        "america-latina", "global", "bilateral",
        
        # Atores
        "governo", "empresa-privada", "ong", "academia",
        "comunidade-local", "indigena", "unfccc"
    ]
}

# ============================================================================
# FUNÇÕES DE CHUNKING (DO SEU CÓDIGO)
# ============================================================================

def chunk_pdf(pdf_path: Path, output_folder: Path, pages_per_chunk: int = 15) -> List[Path]:
    """
    Divide PDF em chunks de N páginas
    Retorna lista de arquivos TXT gerados
    """
    try:
        reader = PdfReader(pdf_path)
        total_pages = len(reader.pages)
        
        # Se for pequeno, não faz chunk
        if total_pages <= CONFIG["chunk_threshold_pages"]:
            text = ""
            for page in reader.pages:
                text += page.extract_text() + "\n"
            
            output_file = output_folder / f"{pdf_path.stem}.txt"
            output_file.write_text(text, encoding='utf-8')
            return [output_file]
        
        # Faz chunking
        print(f"     📊 Documento grande ({total_pages} páginas) - dividindo em chunks")
        
        num_chunks = (total_pages + pages_per_chunk - 1) // pages_per_chunk
        chunk_files = []
        
        for chunk_num in range(num_chunks):
            start_page = chunk_num * pages_per_chunk
            end_page = min(start_page + pages_per_chunk, total_pages)
            
            chunk_text = []
            chunk_text.append(f"DOCUMENTO: {pdf_path.stem}")
            chunk_text.append(f"PARTE {chunk_num + 1} de {num_chunks}")
            chunk_text.append(f"Páginas {start_page + 1} a {end_page}")
            chunk_text.append("=" * 70)
            chunk_text.append("")
            
            for page_num in range(start_page, end_page):
                page = reader.pages[page_num]
                chunk_text.append(f"\n--- PÁGINA {page_num + 1} ---\n")
                chunk_text.append(page.extract_text())
            
            output_file = output_folder / f"{pdf_path.stem}_parte{chunk_num + 1:02d}de{num_chunks:02d}.txt"
            output_file.write_text("\n".join(chunk_text), encoding='utf-8')
            chunk_files.append(output_file)
            
            print(f"        ✅ Parte {chunk_num + 1}/{num_chunks}")
        
        return chunk_files
        
    except Exception as e:
        print(f"     ❌ Erro ao fazer chunk: {e}")
        return []


def chunk_docx(docx_path: Path, output_folder: Path, paragraphs_per_chunk: int = 100) -> List[Path]:
    """
    Divide DOCX em chunks de N parágrafos
    Retorna lista de arquivos TXT gerados
    """
    try:
        doc = Document(docx_path)
        total_paragraphs = len(doc.paragraphs)
        
        # Se for pequeno, não faz chunk
        text = "\n".join([p.text for p in doc.paragraphs])
        
        if len(text) <= CONFIG["chunk_threshold_chars"]:
            output_file = output_folder / f"{docx_path.stem}.txt"
            output_file.write_text(text, encoding='utf-8')
            return [output_file]
        
        # Faz chunking
        print(f"     📊 Documento grande ({total_paragraphs} parágrafos) - dividindo em chunks")
        
        num_chunks = (total_paragraphs + paragraphs_per_chunk - 1) // paragraphs_per_chunk
        chunk_files = []
        
        for chunk_num in range(num_chunks):
            start_para = chunk_num * paragraphs_per_chunk
            end_para = min(start_para + paragraphs_per_chunk, total_paragraphs)
            
            chunk_text = []
            chunk_text.append(f"DOCUMENTO: {docx_path.stem}")
            chunk_text.append(f"PARTE {chunk_num + 1} de {num_chunks}")
            chunk_text.append(f"Parágrafos {start_para + 1} a {end_para}")
            chunk_text.append("=" * 70)
            chunk_text.append("")
            
            for para_num in range(start_para, end_para):
                para_text = doc.paragraphs[para_num].text.strip()
                if para_text:
                    chunk_text.append(para_text)
                    chunk_text.append("")
            
            output_file = output_folder / f"{docx_path.stem}_parte{chunk_num + 1:02d}de{num_chunks:02d}.txt"
            output_file.write_text("\n".join(chunk_text), encoding='utf-8')
            chunk_files.append(output_file)
            
            print(f"        ✅ Parte {chunk_num + 1}/{num_chunks}")
        
        return chunk_files
        
    except Exception as e:
        print(f"     ❌ Erro ao fazer chunk: {e}")
        return []

# ============================================================================
# EXTRAÇÃO DE TEXTO (SIMPLES, SEM CHUNK)
# ============================================================================

def extract_text_simple_pdf(filepath: Path) -> str:
    """Extrai texto de PDF pequeno (sem chunking)"""
    try:
        reader = PdfReader(filepath)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        print(f"  ⚠️ Erro PDF: {e}")
        return ""

def extract_text_simple_docx(filepath: Path) -> str:
    """Extrai texto de DOCX pequeno (sem chunking)"""
    try:
        doc = Document(filepath)
        text = "\n".join([p.text for p in doc.paragraphs])
        return text
    except Exception as e:
        print(f"  ⚠️ Erro DOCX: {e}")
        return ""

def extract_text_from_pptx(filepath: Path) -> str:
    """Extrai texto de PPTX"""
    try:
        prs = Presentation(filepath)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            text += f"\n--- SLIDE {slide_num} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        print(f"  ⚠️ Erro PPTX: {e}")
        return ""

# ============================================================================
# EXTRAÇÃO DE KEYWORDS
# ============================================================================

def _get_anthropic_client() -> anthropic.Anthropic:
    """
    Retorna cliente Anthropic garantindo que uma API key esteja disponível.
    Pode ser definida via CONFIG["anthropic_api_key"] ou variável de ambiente
    ANTHROPIC_API_KEY.
    """
    api_key = CONFIG.get("anthropic_api_key") or os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        raise RuntimeError(
            "ANTHROPIC_API_KEY não definida. Configure CONFIG['anthropic_api_key'] "
            "ou exporte a variável de ambiente antes de usar --keywords-only."
        )
    return anthropic.Anthropic(api_key=api_key)


def extract_keywords_with_claude(text: str, vocabulary: List[str], num_kw: int = 5) -> List[str]:
    """Extrai keywords usando Claude API"""
    try:
        client = _get_anthropic_client()
        
        prompt = f"""Analise o texto e extraia as {num_kw} keywords mais relevantes.

REGRAS:
1. Use APENAS keywords deste vocabulário: {', '.join(vocabulary)}
2. Retorne EXATAMENTE {num_kw} keywords
3. Formato: apenas lista separada por vírgula

TEXTO (primeiros 4000 chars):
{text[:4000]}

Keywords:"""

        message = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=100,
            messages=[{"role": "user", "content": prompt}]
        )
        
        response = message.content[0].text.strip()
        keywords = [kw.strip() for kw in response.split(',')]
        valid = [kw for kw in keywords if kw in vocabulary]
        
        return valid[:num_kw]
        
    except Exception as e:
        print(f"  ⚠️ Erro keywords API: {e}")
        return []

def extract_keywords_local(text: str, vocabulary: List[str], num_kw: int = 5) -> List[str]:
    """Extração local por frequência"""
    text_lower = text.lower()
    
    scores = {}
    for kw in vocabulary:
        search = kw.replace('-', ' ')
        count = text_lower.count(search)
        if count > 0:
            scores[kw] = count
    
    sorted_kw = sorted(scores.items(), key=lambda x: x[1], reverse=True)
    return [kw for kw, _ in sorted_kw[:num_kw]]

# ============================================================================
# PROCESSAMENTO COM CHUNKING INTELIGENTE
# ============================================================================

def process_new_documents():
    """
    Processa novos documentos da pasta input → TXTs em documents
    COM CHUNKING AUTOMÁTICO para documentos grandes
    """
    input_path = Path(CONFIG["input_folder"])
    output_path = Path(CONFIG["documents_folder"])
    
    input_path.mkdir(parents=True, exist_ok=True)
    output_path.mkdir(parents=True, exist_ok=True)
    
    # Busca novos arquivos
    extensions = ['.pdf', '.docx', '.pptx']
    new_files = []
    for ext in extensions:
        new_files.extend(input_path.glob(f"*{ext}"))
    
    if not new_files:
        print("📭 Nenhum arquivo novo para processar")
        return 0
    
    print(f"\n📂 Encontrados {len(new_files)} arquivos novos:")
    
    processed = 0
    total_chunks = 0
    
    for filepath in new_files:
        print(f"\n  📄 {filepath.name}")
        
        # Verifica se já foi processado
        existing_txts = list(output_path.glob(f"{filepath.stem}*.txt"))
        if existing_txts:
            print(f"     ⏭️  Já processado, pulando...")
            continue
        
        print(f"     🔍 Extraindo texto...")
        
        # Processa com ou sem chunking baseado no tipo e tamanho
        chunk_files = []
        
        if filepath.suffix.lower() == '.pdf' and CONFIG["auto_chunk"]:
            # PDF: verifica páginas para decidir se faz chunk
            chunk_files = chunk_pdf(filepath, output_path, CONFIG["pages_per_chunk"])
            
        elif filepath.suffix.lower() == '.docx' and CONFIG["auto_chunk"]:
            # DOCX: verifica caracteres para decidir se faz chunk
            chunk_files = chunk_docx(filepath, output_path, CONFIG["paragraphs_per_chunk"])
            
        elif filepath.suffix.lower() == '.pptx':
            # PPTX: nunca faz chunk (geralmente não são tão grandes)
            text = extract_text_from_pptx(filepath)
            if text and len(text.strip()) >= 50:
                txt_output = output_path / f"{filepath.stem}.txt"
                txt_output.write_text(text, encoding='utf-8')
                chunk_files = [txt_output]
        
        else:
            # Fallback: extração simples sem chunk
            if filepath.suffix.lower() == '.pdf':
                text = extract_text_simple_pdf(filepath)
            elif filepath.suffix.lower() == '.docx':
                text = extract_text_simple_docx(filepath)
            else:
                text = ""
            
            if text and len(text.strip()) >= 50:
                txt_output = output_path / f"{filepath.stem}.txt"
                txt_output.write_text(text, encoding='utf-8')
                chunk_files = [txt_output]
        
        # Contabiliza
        if chunk_files:
            num_files = len(chunk_files)
            if num_files > 1:
                print(f"     ✅ Salvo em {num_files} partes")
                total_chunks += num_files
            else:
                chars = chunk_files[0].read_text(encoding='utf-8')
                print(f"     ✅ Salvo: {len(chars):,} caracteres")
            
            processed += 1
        else:
            print(f"     ❌ Erro no processamento")
    
    if total_chunks > 0:
        print(f"\n📊 {total_chunks} chunks criados de documentos grandes")
    
    return processed

def update_keywords_metadata():
    """
    Atualiza metadata de keywords para todos os TXTs em documents/
    """
    docs_path = Path(CONFIG["documents_folder"])
    metadata_path = Path(CONFIG["metadata_folder"])
    metadata_path.mkdir(parents=True, exist_ok=True)
    
    # Carrega metadata existente
    metadata_file = metadata_path / "keywords_metadata.json"
    if metadata_file.exists():
        with open(metadata_file, 'r', encoding='utf-8') as f:
            all_metadata = json.load(f)
    else:
        all_metadata = {}
    
    # Processa todos os TXTs
    txt_files = list(docs_path.glob("*.txt"))
    
    if not txt_files:
        print("\n⚠️ Nenhum TXT encontrado em documents/")
        return
    
    print(f"\n🏷️  Processando keywords para {len(txt_files)} documentos...")
    
    updated = 0
    for txt_file in txt_files:
        doc_id = txt_file.stem
        
        # Se já tem metadata e keywords, pula
        if doc_id in all_metadata and all_metadata[doc_id].get('keywords'):
            print(f"  ⏭️  {txt_file.name} - já tem keywords")
            continue
        
        print(f"  🔄 {txt_file.name}")
        
        # Lê conteúdo
        text = txt_file.read_text(encoding='utf-8')
        
        # Extrai keywords
        if CONFIG["use_claude_api"]:
            keywords = extract_keywords_with_claude(
                text, 
                CONFIG["keyword_vocabulary"],
                CONFIG["keywords_per_doc"]
            )
        else:
            keywords = extract_keywords_local(
                text,
                CONFIG["keyword_vocabulary"],
                CONFIG["keywords_per_doc"]
            )
        
        if keywords:
            print(f"     🏷️  {', '.join(keywords)}")
        
        # Atualiza metadata
        all_metadata[doc_id] = {
            "filename": txt_file.name,
            "filepath": str(txt_file),
            "keywords": keywords,
            "char_count": len(text),
            "word_count": len(text.split()),
            "processed_date": datetime.now().isoformat()
        }
        
        updated += 1
    
    # Salva metadata
    with open(metadata_file, 'w', encoding='utf-8') as f:
        json.dump(all_metadata, f, indent=2, ensure_ascii=False)
    
    print(f"\n✅ Metadata atualizada: {updated} documentos processados")
    
    # Gera índice invertido
    print("🔄 Gerando índice de keywords...")
    keyword_index = {}
    
    for doc_id, doc_meta in all_metadata.items():
        for kw in doc_meta.get("keywords", []):
            if kw not in keyword_index:
                keyword_index[kw] = []
            keyword_index[kw].append(doc_id)
    
    index_file = metadata_path / "keywords_metadata_index.json"
    with open(index_file, 'w', encoding='utf-8') as f:
        json.dump(keyword_index, f, indent=2, ensure_ascii=False)
    
    print(f"✅ Índice gerado: {len(keyword_index)} keywords únicas")

def sync_with_github():
    """Faz commit e push das mudanças para GitHub"""
    repo_path = Path(CONFIG["github_repo"])
    
    if not (repo_path / ".git").exists():
        print("\n⚠️ Não é um repositório Git")
        return
    
    print("\n📤 Sincronizando com GitHub...")
    
    try:
        os.chdir(repo_path)
        
        # Git add
        subprocess.run(["git", "add", "documents/", "metadata/"], check=True)
        
        # Verifica mudanças
        result = subprocess.run(
            ["git", "status", "--porcelain"],
            capture_output=True,
            text=True
        )
        
        if not result.stdout.strip():
            print("  ℹ️  Nenhuma mudança para commitar")
            return
        
        # Commit
        commit_msg = f"Update knowledge base - {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        subprocess.run(["git", "commit", "-m", commit_msg], check=True)
        print(f"  ✅ Commit: {commit_msg}")
        
        # Push
        if CONFIG["auto_push"]:
            subprocess.run(["git", "push"], check=True)
            print("  ✅ Push realizado")
        else:
            print("  ℹ️  Push desativado")
        
    except subprocess.CalledProcessError as e:
        print(f"  ❌ Erro Git: {e}")
    except Exception as e:
        print(f"  ❌ Erro: {e}")

# ============================================================================
# MAIN
# ============================================================================

def main(keywords_only=False, no_push=False, no_chunk=False):
    """Workflow completo"""
    
    print("\n" + "="*70)
    print("🚀 IETA KNOWLEDGE BASE - WORKFLOW INTEGRADO")
    print("="*70)
    
    if no_push:
        CONFIG["auto_push"] = False
    
    if no_chunk:
        CONFIG["auto_chunk"] = False
        print("⚠️  Chunking automático desabilitado")
    
    # Processa documentos
    if not keywords_only:
        print("\n📄 FASE 1: Processando novos documentos...")
        num_processed = process_new_documents()
        print(f"\n✅ {num_processed} documentos processados")
    else:
        print("\n⏭️  Pulando processamento de documentos")
    
    # Keywords
    if CONFIG["extract_keywords"]:
        print("\n🏷️  FASE 2: Atualizando keywords...")
        update_keywords_metadata()
    
    # GitHub
    if CONFIG["auto_commit"]:
        print("\n📤 FASE 3: Sincronizando com GitHub...")
        sync_with_github()
    
    print("\n" + "="*70)
    print("✨ WORKFLOW CONCLUÍDO!")
    print("="*70)
    print("\n💡 Próximos passos:")
    print("   1. streamlit run meeting_prep.py")
    print("   2. Ou verifique mudanças no GitHub")
    print()

if __name__ == "__main__":
    import argparse
    
    parser = argparse.ArgumentParser(description="Workflow integrado IETA")
    parser.add_argument("--keywords-only", action="store_true",
                       help="Só atualiza keywords")
    parser.add_argument("--no-push", action="store_true",
                       help="Faz commit mas não push")
    parser.add_argument("--no-chunk", action="store_true",
                       help="Desabilita chunking automático")
    parser.add_argument("--local-keywords", action="store_true",
                       help="Usa extração local de keywords")
    
    args = parser.parse_args()
    
    if args.local_keywords:
        CONFIG["use_claude_api"] = False
    
    main(
        keywords_only=args.keywords_only, 
        no_push=args.no_push,
        no_chunk=args.no_chunk
    )